from pptx import Presentation
import numpy as np

# graph layout
from networkx import spectral_layout
# from networkx.drawing.nx_agraph import graphviz_layout

from .utils.coordinate import anchor, get_width_height, CENTER
from .utils.shapes import add_TextBox, add_connector

LAYOUT = spectral_layout

# color constants
GRAY  = (210, 210, 210)
GREEN = (195, 214, 155)
BLACK = (  0,   0,   0)


def fit_slide_shape(pos):
    """
    Min-max scaling to slide shape
    """
    points = np.array(list(pos.values()))  # shape = (n, 2)
    min_xy = points.min(axis=0)
    max_xy = points.max(axis=0)
    scale_xy = max_xy - min_xy
    ratio_xy = CENTER / scale_xy * 2
    for name, p in pos.items():
        pos[name] = ((p - min_xy) * ratio_xy).round().astype(int)


def scale_up(pos, scale=1, origin="center"):
    """
    Notice:
        input origin = `origin`
        scaling center = CENTER
        output origin = upper-left corner

    Parameters
    ----------
    pos : dict
        { node_name: node_position ... }
    scale : float or np.array
        scale parameter
    origin : {'center', 'upper-left'}
        Origin of input position expression
    """
    for name, p in pos.items():
        if origin == "center":
            pos[name] = (scale * p + CENTER).round().astype(int)
        else:
            pos[name] = (scale * (p - CENTER) + CENTER).round().astype(int)


def convert(graph, path, layout=LAYOUT, scale=1, groups="files", connector_type=1):
    """
    Convert a networkx DiGraph into PowerPoint slide object

    Parameters
    ----------
    graph : networkx.Digraph
        Set following node attributes :
            'content' str, default " "
            'color' (int, int, int), default (195, 214, 155)
                rgb color of title box
            'bg' (int, int, int), default (210, 210, 210)
                rgb color of content box
    path : str
        ex: 'ooxx/abc/hello_world.pptx'
    layout : {func, dict}, default networkx.circular_layout
        Graph layout function or dict as { node: [x, y] ... },
        provide relative position of nodes.
    scale : float or np.array, default 1
        If it is float array of size 2, it means scale size of x and y.
        If it is float, both x and y are scaled as it.
        Negative means reflection on the axis.
    groups : {'colors', 'files'}, default 'colors'
        'colors' group title boxes together and content boxes together
        'files' group title box and content box of each file
    connector_type : {1, 2, 3}, default 1
        1 = STRAIGHT
        2 = ELBOW
        3 = CURVE
    """
    G = graph.copy()
    file = Presentation()
    page = file.slides.add_slide(file.slide_layouts[6])  # blank layout

    # positions
    pos = layout if type(layout) == dict else layout(G)
    fit_slide_shape(pos)
    scale_up(pos, scale=scale, origin="upper-left")

    for node, (px, py) in pos.items():
        attr = G.nodes[node]
        title = str(node)
        content = str(attr.get("content", " "))
        # calculate width & anchor
        title_width, title_height = get_width_height(title)
        content_width, content_height = get_width_height(content)

        width = max(title_width, content_width)
        height = title_height + content_height
        x, y = anchor(px, py, width, height)

        # textbox
        attr["title"] = add_TextBox(
            page, title, x, y, width, attr.get("color", GREEN)
        )
        attr["content"] = add_TextBox(
            page, content, x, y + title_height, width, attr.get("bg", GRAY)
        )

    link_shapes = [
        add_connector(
            page,
            G.nodes[u]["title"],
            G.nodes[u]["content"],
            G.nodes[v]["title"],
            G.nodes[v]["content"],
            connector_type,
        )
        for u, v in G.edges
    ]
    page.shapes.add_group_shape(link_shapes)  # group all connector together

    if groups == "files":
        for node in G.nodes:
            attr = G.nodes[node]
            page.shapes.add_group_shape([attr["title"], attr["content"]])
    else:  # groups == 'colors'
        # group all title boxes
        page.shapes.add_group_shape(
            [box for n, box in G.nodes.data("title")]
        )
        # group all content boxes
        page.shapes.add_group_shape(
            [box for n, box in G.nodes.data("content")]
        )

    file.save(path)
