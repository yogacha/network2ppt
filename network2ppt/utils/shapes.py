import numpy as np
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

from .coordinate import get_width_height


def add_TextBox(slide, text, x, y, width=-1, rgb=None):
    """
    create a textbox automatically adjust height by it's content
    font--consolas

    Parameters
    ----------
    slide :
        slide to add textbox
    text : str
        content string
    x : int
        position on x-axis
    y : int
        position on y-axis
    width : int, default -1
        If -1, set to be minimum width that can fit the content
    rgb : (int, int, int), default None
        If None, set to be transparent
    """

    # width
    min_width, height = get_width_height(text)
    width = max(width, min_width)

    # textbox
    box = slide.shapes.add_textbox(x, y, width, height)
    box.text = text
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    box.text_frame.paragraphs[0].font.name = 'Consolas'
    
    # rgb color
    if rgb:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*rgb)

    return box


def add_connector(page, t1, c1, t2, c2, connector_type):
    """
    Adding shortest connector between (title1, content1) & (title2, content2)

     +------ 0 ------+
     |               |
     1   Title box   3
     |               |
     +------ 2 ------+
     +------ 0 ------+
     |               |
     1  Content box  3
     |               |
     +------ 2 ------+
    """
    line = page.shapes.add_connector(connector_type, 0, 0, 0, 0)
    pair1 = [(t1, 0), (t1, 1), (c1, 1), (c1, 2), (c1, 3), (t1, 3)]
    pair2 = [(t2, 0), (t2, 1), (c2, 1), (c2, 2), (c2, 3), (t2, 3)]

    min_dist = np.inf
    for box1, i in pair1:
        for box2, j in pair2:
            line.begin_connect(box1, i)
            line.end_connect(box2, j)
            dist = (line.begin_x - line.end_x) ** 2 + (line.begin_y - line.end_y) ** 2
            if dist < min_dist:
                min_dist = dist
                res = box1, i, box2, j
    box1, i, box2, j = res
    line.begin_connect(box1, i)
    line.end_connect(box2, j)

    return line
